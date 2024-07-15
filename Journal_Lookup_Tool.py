
"""
v3.0.0:

7/10/24
Lisa P.

Refactoring Aaron's Journal_Lookup_Tool

adapt entrez query error handling from 
https://stackoverflow.com/questions/75678873/unpredictable-httperror-using-entrez-esearch-function-from-bio-package

TODO: improve sql query name sensitivity
TODO: refactor get_pptx

7/10/24
APL

- re-incorporated sleep(1) between requests, else XML errors from Entrez
- improved formatting for google slies by reducing title character limit
- if using current date, ppt now displays current date instead of '3000'

"""


####################################################################################################
# IMPORTS
####################################################################################################
from Bio import Entrez
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import time
#import requests
from pptx.dml.color import RGBColor
from urllib.error import HTTPError



####################################################################################################
# GLOBAL VARIABLES
####################################################################################################
version = "3.0.0"
update_date_numbers = 20240710 # last update



####################################################################################################
# FUNCTIONS
####################################################################################################

##############################
# A. Error handler, getters
##############################
# function to check user-input date formatting:
def remove_duplicates(components):
    seen = set()
    return [x for x in components if not (x["Title"] in seen or seen.add(x["Title"]))]

def validate_date(date):
    try:
        datetime.strptime(date, "%Y/%m/%d")
    except ValueError:
        raise ValueError("Incorrect date format, should be YYYY/MM/DD")

def get_pmids(combined_query: str):
    """
    Take in a sql query and search pubmed using Entrez.esearch
    """
    try:
        handle = Entrez.esearch(db="pubmed", 
                                term=combined_query, 
                                retmax=100000)
        try:    
            record = Entrez.read(handle, validate=False)
        except:
            raise ValueError(f"Entrez.read error. Error reading XML for query {combined_query}")
        #print(record)
        pmids = record["IdList"]
        return pmids

    except HTTPError as e:
        raise HTTPError(f"Entrez.esearch error. Error retrieving XML for query {combined_query}: {e}")

def get_paper(pmid: str):
    """
    Use entrez.efetch to get paper metadata
    
    """
    
    try:
        handle = Entrez.efetch(db="pubmed", 
                               id=pmid, 
                               rettype="medline", 
                               retmode="xml")
        paper = Entrez.read(handle, validate=False)
        #print(f"Paper with PMID {pmid} found and added")
        #print("\n")
    except HTTPError as e:
        if (e.response.status_code == 500):  # If the error is an HTTP 500 error, retry
            print("HTTP Error 500 encountered. Retrying...")
            time.sleep(10)  # Wait for 10 second before retrying
        elif (e.response.status_code == 429):  # If the error is an HTTP 429 error, retry after 10 seconds
            print("HTTP Error 429 encountered. Too many requests. Retrying after 10 seconds...")
            time.sleep(10)  # Wait for 10 seconds before retrying
        else:
            raise e# If the error is not an HTTP 500 or 429 error, raise it
    return paper

def get_paper_info(papers: list) -> list:
    """
    Extract fields in the paper relevant for journal club presentation
    """
    # Update the components list based on the updated papers list
    components = []  
    for paper in papers:
        for article in paper["PubmedArticle"]:
            if "MedlineCitation" in article and "Article" in article["MedlineCitation"]:
                component = {}
                # Extract the title                
                try:
                    component["Title"] = article["MedlineCitation"]["Article"]["ArticleTitle"]
                except KeyError:
                    component["Title"] = "No title available"
                    #break out of for loop if no title available
                    break
                
                # Extract the abstract
                try:
                    component["Abstract"] = article["MedlineCitation"]["Article"]["Abstract"]["AbstractText"][0]
                except KeyError:
                    component["Abstract"] = "No abstract available"
                    break

                # Extract the journal name
                if (
                    "Journal" in article["MedlineCitation"]["Article"]
                    and "Title" in article["MedlineCitation"]["Article"]["Journal"]
                ):
                    component["Journal"] = article["MedlineCitation"]["Article"]["Journal"]["Title"]

                # Extract the publication date
                try:
                    component["Date"] = article["MedlineCitation"]["Article"]["ArticleDate"]
                except KeyError:
                    component["Date"] = "No date available"
                
                # Extract the link to the publication
                try:
                    component["Link"] = article["MedlineCitation"]["Article"]["ELocationID"]
                except KeyError:
                    component["Link"] = "No link available"

                # Extract the authors
                try:
                    component["Authors"] = [
                        f"{author['LastName']}"
                        for author in article["MedlineCitation"]["Article"]["AuthorList"]
                    ]
                except KeyError:
                    component["Authors"] = "No authors available"
                
                # Extract the affiliations
                if "AuthorList" in article["MedlineCitation"]["Article"]:
                    component["Institution"] = [
                        str(affiliation["Affiliation"])
                        for author in article["MedlineCitation"]["Article"]["AuthorList"]
                        if "AffiliationInfo" in author
                        for affiliation in author["AffiliationInfo"]
                    ]
                components.append(component)

    components = remove_duplicates(components)
    return components

def get_pptx(start_end_date: tuple,
             config_file_dict: dict, 
             components: list, 
             pptx_name: str = 'publications.pptx', 
             version: str = version, 
             update_date_numbers: int = update_date_numbers, ):
    """
    Producing the powerpoint file:
        - Creates an intro / about slide
        - Creates a summary of articles found slide
        - Creates a slide-per-abstract and relevant information, ready for import into google sheets website!
            - this can be copy/pasted slide-wise into the journal club Google Slides Presentation
    """
    #check to not overwrite existing file
    if os.path.exists(pptx_name):
        print(f"{pptx_name} already exists. Please rename or delete it.")
        return
        
    # Create a new presentation in powerpoint
    presentation = Presentation()

    # Retrieve the start and end dates
    start_date, end_date = start_end_date

    # Retrieve the journals of interest
    email = config_file_dict['email']
    journals_of_interest = config_file_dict['journals']
    topics = config_file_dict['topics']

    # slide for the introduction
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"Journal Lookup Tool v{version}"
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"By AP Ledray, updated {update_date_numbers}"
    p.space_after = Pt(14)
    p = tf.add_paragraph()
    p.text = "contact: aaronledray@gmail.com"
    p = tf.add_paragraph()
    p.text = ""

    # slide for the username, dates, and journals
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Query Information (1):"
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"Username: {email}"
    run.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = f"Query start date: {start_date}\n"
    if end_date == "3000/01/01":
        current_date = datetime.now()
        current_formatted_date = current_date.strftime("%Y/%m/%d")
        run.text += f"Query end date: {current_formatted_date}\n"
    else:
        run.text += f"Query end date: {end_date}\n"
    run.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Query journals:"
    run.font.size = Pt(16)
    for journal in journals_of_interest:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = journal
        run.font.size = Pt(12)

    # slide for the keywords
    max_characters_per_line = 100
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Query Information (2):"
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Query keywords:"
    p = tf.add_paragraph()
    run.font.size = Pt(16)
    # Need to first turn all keywrods of interest into a single string separated by commas.
    combined_keyword_queries = ", ".join(topics)
    if len(combined_keyword_queries) > max_characters_per_line:
        while len(combined_keyword_queries) > max_characters_per_line:
            # Find the last space before the Xth character
            split_index = combined_keyword_queries.rfind(" ", 0, max_characters_per_line)
            if split_index == -1:
                split_index = max_characters_per_line
            chunk = combined_keyword_queries[:split_index]
            combined_keyword_queries = combined_keyword_queries[
                split_index:
            ].lstrip()  # Remove leading spaces from the remaining text
            run = p.add_run()
            run.text = chunk + "\n"
            # print (chunk)
            run.font.size = Pt(12)
        # Add the remaining text
        run = p.add_run()
        run.text = combined_keyword_queries
        # print (title_plain_text)
        run.font.size = Pt(12)
    else:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = combined_keyword_queries
        run.font.size = Pt(12)

    # Initialize a dictionary to store the article counts for each journal
    # Journal_data is a weird name for it, but it's the frequency of each journal hit
    journal_data = {}
    for component in components:
        if "Journal" in component:
            journal_name = component["Journal"]
            journal_data[journal_name] = journal_data.get(journal_name, 0) + 1

    #slide for the journal summary table
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Journals Found:"
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)
    table = slide.shapes.add_table(
        rows=len(journal_data) + 1,
        cols=2,
        left=left,
        top=top,
        width=width,
        height=height,
    ).table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    table.cell(0, 0).text = "Journal"
    table.cell(0, 1).text = "Number of Articles"
    black = RGBColor(0, 0, 0)
    # Defining colors for the table:
    light_grey = RGBColor(211, 211, 211)
    dark_grey = RGBColor(169, 169, 169)
    # Set the color of the header row to black
    for cell in table.rows[0].cells:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = black
    # Populate the table with journal data
    for i, (journal, num_articles) in enumerate(journal_data.items()):
        row = table.rows[i + 1]
        row.cells[0].text = journal
        row.cells[1].text = str(num_articles)

        # Alternate row colors for both cells
        fill_0 = row.cells[0].fill
        fill_1 = row.cells[1].fill
        fill_0.solid()
        fill_1.solid()

        if i % 2 == 0:
            fill_0.fore_color.rgb = light_grey
            fill_1.fore_color.rgb = light_grey
        else:
            fill_0.fore_color.rgb = dark_grey
            fill_1.fore_color.rgb = dark_grey 
        
    # Paper-wise, so it's for all i in components:
    for component in components:
        # Add a slide with a blank layout
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)
        left = top = Inches(0)  # Sets the left and top position of the text box to 1 inch from the slide's left edge            
        width = height =  presentation.slide_width - Inches(0.5)  # Adjust the width to fit the slide and set equal to height
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        # Add the article's title text
        p = tf.add_paragraph()
        run = p.add_run() 
        max_characters_per_line = 75 # adjusted to 75 from 80
        # TODO: split this replace chain into a function so you can call it in both spots
        title_plain_text = (
            component["Title"]
            .replace("<sub>", "")
            .replace("</sub>", "")
            .replace("<i>", "")
            .replace("</i>", "")
            .replace("<sup>", "")
            .replace("</sup>", "")
        )
        if len(title_plain_text) > max_characters_per_line:
            while len(title_plain_text) > max_characters_per_line:
                # Find the last space before the Xth character
                split_index = title_plain_text.rfind(" ", 0, max_characters_per_line)
                if split_index == -1:
                    split_index = max_characters_per_line
                chunk = title_plain_text[:split_index]
                title_plain_text = title_plain_text[
                    split_index:
                ].lstrip()  # Remove leading spaces from the remaining text
                run = p.add_run()
                run.text = chunk + "\n"
                # print (chunk)
                run.font.size = Pt(20)
            # Add the remaining text
            run = p.add_run()
            run.text = title_plain_text
            # print (title_plain_text)
            run.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = title_plain_text
            run.font.size = Pt(20)
        # Add a blank line
        tf.add_paragraph()
        # Add the abstract text
        max_characters_per_line = 115
        p = tf.add_paragraph()
        run = p.add_run()
        if "Abstract" in component:
            abstract_plain_text = (
                component["Abstract"]
                .replace("<sub>", "")
                .replace("</sub>", "")
                .replace("<i>", "")
                .replace("</i>", "")
                .replace("<sup>", "")
                .replace("</sup>", "")
            )
            if len(abstract_plain_text) > max_characters_per_line:
                while len(abstract_plain_text) > max_characters_per_line:
                    # Find the last space before the Xth character
                    split_index = abstract_plain_text.rfind(
                        " ", 0, max_characters_per_line
                    )
                    if split_index == -1:
                        # If no space was found, split at the Xth character
                        split_index = max_characters_per_line
                    chunk = abstract_plain_text[:split_index]
                    abstract_plain_text = abstract_plain_text[
                        split_index:
                    ].lstrip()  # Remove leading spaces from the remaining text
                    run = p.add_run()
                    run.text = chunk + "\n"  # Add a newline at the end of each chunk
                    run.font.size = Pt(14)
                # Add the remaining text
                run = p.add_run()
                run.text = abstract_plain_text
                run.font.size = Pt(14)
            else:
                run = p.add_run()
                run.text = abstract_plain_text
                run.font.size = Pt(14)
        else:
            run.text = "Abstract: Not available"
            run.font.size = Pt(14)
        # Add a blank line
        tf.add_paragraph()
        # Add the journal name and publication date
        p = tf.add_paragraph()
        run = p.add_run()
        # Check if 'journal' is in the component
        if "Journal" in component:
            run.text += "Published in: " + component["Journal"] + "\n"
            run.font.size = Pt(12)
        else:
            run.text += "Journal: Not available\n"
            run.font.size = Pt(12)

        # Check if 'date' is in the component
        if "Date" in component:
            # print ("Date found!")
            date = component["Date"][0]  # Get the first date
            # Check if 'Year', 'Month', and 'Day' are in the date
            if "Year" in date and "Month" in date and "Day" in date:
                date_str = "{}/{}/{}".format(
                    date["Year"], date["Month"], date["Day"]
                )  # Format the date as a string
                run.text += "Publication Date: " + date_str + "\n"
        else:
            run.text += "Publication Date: not available\n"
        # Add the DOI
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "DOI: " + component["Link"][0] + "\n"
        run.font.size = Pt(10)  # Set the font size
        # Add the Authors
        max_characters_per_line = 170
        p = tf.add_paragraph()
        run = p.add_run()
        if "Authors" in component:
            authors_plain_text = "Authors: " + ", ".join(component["Authors"])
            if len(authors_plain_text) > max_characters_per_line:
                while len(authors_plain_text) > max_characters_per_line:
                    # Find the last space before the Xth character
                    split_index = authors_plain_text.rfind(
                        " ", 0, max_characters_per_line
                    )
                    if split_index == -1:
                        # If no space was found, split at the Xth character
                        split_index = max_characters_per_line
                    chunk = authors_plain_text[:split_index]
                    authors_plain_text = authors_plain_text[
                        split_index:
                    ].lstrip()  # Remove leading spaces from the remaining text
                    run = p.add_run()
                    run.text = chunk + "\n"
                    run.font.size = Pt(10)  # Set the font size
                # Add the remaining text
                run = p.add_run()
                run.text = authors_plain_text
                run.font.size = Pt(10)  # Set the font size
            else:
                # p = tf.add_paragraph()
                run = p.add_run()
                run.text = authors_plain_text
                run.font.size = Pt(10)  # Set the font size
        else:
            run.text = "Authors: Not available"
            run.font.size = Pt(10)  # Set the font size
        # Add a blank line
        tf.add_paragraph()
        # Add the Institutions
        p = tf.add_paragraph()
        run = p.add_run()
        max_characters_per_line = 170
        if "Institution" in component:
            # Count the frequency of each institution
            institution_count = {}
            for institution in component["Institution"]:
                if institution in institution_count:
                    institution_count[institution] += 1
                else:
                    institution_count[institution] = 1


        # Sort the institutions by frequency and select the top two    
        top_two_institutions = sorted(institution_count, key=institution_count.get, reverse=True)[:2]
        # Create a list of strings containing the institution names and their counts, making combined string
        institution_info = [f"{institution} {institution_count[institution]}" for institution in top_two_institutions]
        # Joining the list of strings into a single string
        institution_text = "Institutions: " + "; ".join(institution_info)
        # Check if the text is too long to print directly and needs to be formatted:
        if len(institution_text) > max_characters_per_line:
            # Split the text into chunks and add newlines
            while len(institution_text) > max_characters_per_line:
                split_index = institution_text.rfind(" ", 0, max_characters_per_line)
                if split_index == -1:
                    split_index = max_characters_per_line
                chunk = institution_text[:split_index]
                institution_text = institution_text[split_index:].lstrip()
                run = p.add_run()
                run.text = chunk + "\n"
                run.font.size = Pt(10)
            # Add the remaining text
            run = p.add_run()
            run.text = institution_text
            run.font.size = Pt(10)
        else:
            run.text = institution_text
            run.font.size = Pt(10)

    # Finally, save the presentation!
    presentation.save("publications.pptx")


#############################
# B. Reader, searcher, writer
#############################

def read_config_file(config_file : str = 'config.txt') -> dict:
    """
    config.txt is a file with three paragraph
    first paragraph is user email
    second paragraph is journal name, one per line
    third paragraph is topic, one per line

    Read config file and return dictionary of
    {'email': user_email, 
    'journals': [journal1, journal2, ...],
    'topics': [topic1, topic2, ...]}

    """
    if os.path.exists(config_file):
        print("config.txt found!")
        #time.sleep(1)
        #print("\n")
        with open(config_file, "r") as f:
            lines = f.read()
            #check for empty file
            if len(lines) == 0:
                raise ValueError("config.txt is empty")
            #check for exactly 3 paragraph
            try:
                email, journals, topic = lines.split('\n\n')
            except:
                raise ValueError('config.txt improperly formatted. ' + 
                      'Please make sure there are three paragraphs separated by two newlines')
            
            #make sure only one email is provided
            if email.count('@') != 1:
                raise ValueError("Please provide only one email")

            print(f"Email: {email}")
            print("\n")
            print(f"Journals: \n{journals}")
            print("\n")
            print(f"Topics: \n{topic}")

        #create a dictionary
        config_file_dict = {}
        config_file_dict['email'] = email.strip()
        config_file_dict['journals'] = journals.splitlines()
        config_file_dict['topics'] = topic.splitlines()

        return config_file_dict

def ask_user_date(end_date_default: str = '3000/01/01') -> tuple:
    """
    Ask user for time span to search
    Default: 1 week ago until current date
    Max year is 3000
    Return tuple of start_date, end_date
    """
    # Check if user wants default settings, or different timespan:
    print("\n")
    running_mode = input(
        "Would you like to run in default mode? " + 
        "Default mode is query from 1 week ago until current date. (y)es/(n)o: "
    )
    if running_mode.lower() in ("y", "yes", ""):
        end_date = datetime.now()
        start_date = end_date - timedelta(weeks=1)
        start_date = start_date.strftime("%Y/%m/%d")
        end_date = end_date_default  # Why is 3000? Basically it means that the maximum date queried is the year 3000.
    else:
        start_date = input(
            "Please enter the start date (YYYY/MM/DD), for example 2020/01/01: "
        )
        validate_date(start_date)
        end_date = input("Please enter the end date (YYYY/MM/DD), for example 3000/01/01: ")
        if end_date == "":
            end_date = end_date_default
        else:
            validate_date(end_date)

        start_date = datetime.strptime(start_date, "%Y/%m/%d")
        end_date = datetime.strptime(end_date, "%Y/%m/%d")
        

    print(f"Start date: {start_date}")
    print(f"End date: {end_date}")    

    return (start_date, end_date)

def lookup_pubmed(config_file_dict : dict, start_end_date: tuple, attempt_number: int = 2) -> list:
    """
    Use the config file dictionary to search pubmed for articles
    Pubmed is queried with keyword combos

    Sql query of Pubmed via Entrez.esearch using string of the format:
    '{keyword_query} AND ("{start_date}"[Date - Entry] : "{end_date}"[Date - Entry]) AND "{journal_of_interest}"[Journal]'

    return a list of papers
    """
    # Initialize the lists for papers and components
    papers = []

    #retrieve values from config file
    email = config_file_dict['email']
    journals = config_file_dict['journals']
    keywords = config_file_dict['topics']

    #retrieve values from start_end_date tuple
    start_date, end_date = start_end_date

    # Set the email for Entrez
    Entrez.email = email

    # Fetch PMIDs based on the keyword query
    for journal_of_interest in journals:
        print(f'Searching journal: {journal_of_interest}')
        for keyword_query in keywords:
            # Construct the combined query for pubmed:
            combined_query = f'{keyword_query} AND ("{start_date}"[Date - Entry] : "{end_date}"[Date - Entry]) AND "{journal_of_interest}"[Journal]'
            #print(f"Pubmed query keywords: {combined_query}")
            
            # Fetch the PMIDs for the query
            pmids = get_pmids(combined_query)
            #print(f"PMIDs found: {pmids}")
            
            # for successful PMIDs found from the queries:
            # Fetch the metadata of the papers for each PMID
            for pmid in pmids:
                # THIS IS IMPORTANT, or requests will be denied:
                time.sleep(1)  # Add a delay of 1 second between each request    
                for _ in range(attempt_number):  # Try (x) times!
                    try:
                        paper = get_paper(pmid)
                        papers.append(paper)
                    except HTTPError as e:
                        print(f"HTTPError encountered: {e}")
                        print(f"Failed to retrieve paper with PMID {pmid}")
                        continue
    return papers



####################################################################################################
# C. PRINTERS
####################################################################################################



def print_opener():
    print(rf"""
--------------------------------------------------

Journal Lookup Tool by AP Ledray

Version: {version}

updated: {update_date_numbers}
 ____ ____ ____ ____ ____ ____ ____
||J |||o |||u |||r |||n |||a |||l ||
||__|||__|||__|||__|||__|||__|||__||
|/__\|/__\|/__\|/__\|/__\|/__\|/__\|
 ____ ____ ____ ____
||C |||l |||u |||b ||
||__|||__|||__|||__||
|/__\|/__\|/__\|/__\|
 ____ ____ ____ ____ ____ ____
||L |||o |||o |||k |||u |||p ||
||__|||__|||__|||__|||__|||__||
|/__\|/__\|/__\|/__\|/__\|/__\|
 ____ ____ ____ ____
||T |||o |||o |||l ||
||__|||__|||__|||__||
|/__\|/__\|/__\|/__\|

--------------------------------------------------)
""")

# function for closer printing 1
def print_closer_0():
     print(r"""
--------------------------------------------------

  /^ ^\
 / 0 0 \    Oh, heck!
 V\ Y /V
  / - \
 /    |
V__) ||

No papers were found using query information, publications.pptx has not been updated.
        
--------------------------------------------------
    """)

# function for closer printing 2
def print_closer_1():
    print(r"""
--------------------------------------------------

     .-'
'--./ /     _.---.
'-,  (__..-`       \
   \          .     |    Aw yea!
    `,.__.   ,__.--/
      '._/_.'___.-`

Journal_Lookup_Tool has finished, publications.pptx has been updated!

Have a nice day!
        
--------------------------------------------------
    """)



####################################################################################################
# MAIN
####################################################################################################




def main():
    print_opener()
    #1. Read config file
    print('Reading config file...   ')
    config_file_dict = read_config_file('config.txt')
    
    #2. Get date range
    start_end_date = ask_user_date()
    print('Looking up papers...   ')
    papers = lookup_pubmed(config_file_dict=config_file_dict, 
                start_end_date=start_end_date)
    
    #3. Get paper info
    print('Getting paper info...   ')
    components = get_paper_info(papers)
    #print(components)

    if len(components) == 0:
        print_closer_0()
    else:
        print_closer_1()

    #4. Prepare pptx
    print('Prepare pptx... ')
    get_pptx(start_end_date = start_end_date,
             config_file_dict = config_file_dict,
             components = components)
    
if __name__ == "__main__":
    main()